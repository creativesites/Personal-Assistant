"""
Knowledge Base Retrieval — Phase 8.

Processes uploaded documents (chunking + embedding) and retrieves relevant chunks
for agent context injection during autonomous response generation.
"""

import re
import csv
import io
from pathlib import Path
import structlog
import httpx
import numpy as np
import pdfplumber
from openpyxl import load_workbook
from ..ai.client import get_ai_client
from ..config import settings
from ..database import get_pool

log = structlog.get_logger()

# Target chunk size in words; hard cap prevents over-large chunks
_TARGET_CHUNK_WORDS = 500
_MAX_CHUNK_WORDS = 600

# Simple HTML tag pattern used for stripping markup from fetched URLs
_HTML_TAG_RE = re.compile(r'<[^>]+>')
_WHITESPACE_RE = re.compile(r'\s{2,}')


def _strip_html(html: str) -> str:
    """Strip HTML tags and normalise whitespace, returning plain text."""
    text = _HTML_TAG_RE.sub(' ', html)
    text = _WHITESPACE_RE.sub(' ', text)
    return text.strip()


def _split_into_chunks(text: str) -> list[str]:
    """
    Split text into chunks of roughly _TARGET_CHUNK_WORDS words.

    Splits on paragraph boundaries first, then accumulates paragraphs until
    the target word count is reached. A paragraph that alone exceeds
    _MAX_CHUNK_WORDS is hard-split at sentence boundaries.
    """
    paragraphs = [p.strip() for p in re.split(r'\n{2,}|\r\n\r\n', text) if p.strip()]

    chunks: list[str] = []
    current_words: list[str] = []
    current_count = 0

    for para in paragraphs:
        para_words = para.split()
        para_count = len(para_words)

        if para_count > _MAX_CHUNK_WORDS:
            # Flush what we have first
            if current_words:
                chunks.append(' '.join(current_words))
                current_words = []
                current_count = 0
            # Hard-split on sentence boundaries
            sentences = re.split(r'(?<=[.!?])\s+', para)
            sentence_buf: list[str] = []
            sentence_count = 0
            for sentence in sentences:
                s_words = sentence.split()
                if sentence_count + len(s_words) > _TARGET_CHUNK_WORDS and sentence_buf:
                    chunks.append(' '.join(sentence_buf))
                    sentence_buf = s_words
                    sentence_count = len(s_words)
                else:
                    sentence_buf.extend(s_words)
                    sentence_count += len(s_words)
            if sentence_buf:
                chunks.append(' '.join(sentence_buf))
            continue

        if current_count + para_count > _TARGET_CHUNK_WORDS and current_words:
            chunks.append(' '.join(current_words))
            current_words = para_words
            current_count = para_count
        else:
            current_words.extend(para_words)
            current_count += para_count

    if current_words:
        chunks.append(' '.join(current_words))

    return [c for c in chunks if c]


def _estimate_tokens(text: str) -> int:
    """Rough token estimate: ~0.75 words per token (4 chars avg)."""
    return max(1, len(text) // 4)


def _safe_storage_path(storage_path: str) -> Path:
    path = Path(storage_path)
    if not path.is_absolute():
        path = Path(settings.kb_storage_dir) / path
    resolved = path.resolve()
    root = Path(settings.kb_storage_dir).resolve()
    if root not in resolved.parents and resolved != root:
        raise ValueError('Document storage path is outside KB storage directory')
    return resolved


def _extract_pdf_text(file_path: Path) -> str:
    pages: list[str] = []
    with pdfplumber.open(str(file_path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ''
            if text.strip():
                pages.append(f'=== Page {idx} ===\n{text.strip()}')
    return '\n\n'.join(pages)


def _extract_workbook_text(file_path: Path) -> str:
    workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(' | '.join(values))
        if rows:
            sections.append(f"=== Sheet: {sheet.title} ===\n" + '\n'.join(rows))
    workbook.close()
    return '\n\n'.join(sections)


def _extract_csv_text(file_path: Path) -> str:
    raw = file_path.read_text(encoding='utf-8', errors='ignore')
    rows = csv.reader(io.StringIO(raw))
    return '\n'.join(' | '.join(cell.strip() for cell in row if cell.strip()) for row in rows)


def _extract_docx_text(file_path: Path) -> str:
    try:
        import docx
        doc = docx.Document(str(file_path))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return '\n\n'.join(paras)
    except Exception:
        raw = file_path.read_text(encoding='utf-8', errors='ignore')
        return _strip_html(raw)


def _extract_pptx_text(file_path: Path) -> str:
    try:
        import pptx
        prs = pptx.Presentation(str(file_path))
        slides: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"=== Slide {idx} ===\n" + '\n'.join(texts))
        return '\n\n'.join(slides)
    except Exception:
        raw = file_path.read_text(encoding='utf-8', errors='ignore')
        return _strip_html(raw)


async def _extract_file_text(
    source_type: str, storage_path: str, mime_type: str | None, user_id: str | None = None,
) -> str:
    file_path = _safe_storage_path(storage_path)
    if not file_path.exists():
        raise ValueError('Uploaded file is missing from KB storage')

    st = source_type.lower()
    if st == 'pdf':
        return _extract_pdf_text(file_path)
    if st in ('excel', 'xlsx', 'xls'):
        return _extract_workbook_text(file_path)
    if st == 'csv':
        return _extract_csv_text(file_path)
    if st in ('word', 'docx', 'doc'):
        return _extract_docx_text(file_path)
    if st in ('pptx', 'ppt', 'presentation'):
        return _extract_pptx_text(file_path)
    if st == 'image':
        client = get_ai_client()
        return await client.extract_image_text(
            image_bytes=file_path.read_bytes(),
            mime_type=mime_type or 'image/jpeg',
            service='intelligence', feature='ocr_extraction', user_id=user_id,
        )
    return file_path.read_text(encoding='utf-8', errors='ignore')



async def process_document(document_id: str, user_id: str) -> None:
    """
    Process a KB document by chunking its content and generating embeddings.

    Steps:
      1. Fetch the kb_documents row.
      2. If source_type == 'url', fetch and strip the page HTML.
      3. Split the raw_content into word-count-bounded chunks.
      4. Embed each chunk via the AI client.
      5. Insert kb_chunks rows.
      6. Update kb_documents.status = 'ready' and chunk_count = N.
      On any error, set status = 'error' and record the message.
    """
    pool = await get_pool()

    async with pool.acquire() as conn:
        doc = await conn.fetchrow(
            """
            SELECT id, user_id, agent_id, title, source_type, source_url,
                   raw_content, storage_path, mime_type, status
            FROM kb_documents
            WHERE id = $1 AND user_id = $2
            """,
            document_id,
            user_id,
        )

    if not doc:
        log.warning('kb_document_not_found', document_id=document_id)
        return

    log.info(
        'kb_document_processing_start',
        document_id=document_id,
        source_type=doc['source_type'],
        title=doc['title'],
    )

    try:
        raw_content: str = doc['raw_content'] or ''

        # Fetch URL content if needed
        if doc['source_type'] == 'url' and doc['source_url']:
            async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as http:
                resp = await http.get(doc['source_url'])
                resp.raise_for_status()
                raw_content = _strip_html(resp.text)
            # Persist fetched content back to the document
            async with pool.acquire() as conn:
                await conn.execute(
                    'UPDATE kb_documents SET raw_content = $1 WHERE id = $2',
                    raw_content,
                    document_id,
                )

        if not raw_content.strip() and doc['storage_path']:
            raw_content = await _extract_file_text(
                doc['source_type'],
                doc['storage_path'],
                doc['mime_type'],
                user_id=user_id,
            )
            word_count = len(raw_content.split())
            async with pool.acquire() as conn:
                await conn.execute(
                    """
                    UPDATE kb_documents
                    SET raw_content = $1, word_count = $2, extracted_at = NOW(), updated_at = NOW()
                    WHERE id = $3
                    """,
                    raw_content,
                    word_count,
                    document_id,
                )

        if not raw_content.strip():
            raise ValueError('Document has no text content to process')

        async with pool.acquire() as conn:
            await conn.execute('DELETE FROM kb_chunks WHERE document_id = $1', document_id)

        chunks = _split_into_chunks(raw_content)
        if not chunks:
            raise ValueError('Document produced no chunks after splitting')

        log.info('kb_document_chunked', document_id=document_id, chunk_count=len(chunks))

        client = get_ai_client()
        inserted = 0

        async with pool.acquire() as conn:
            for idx, chunk_text in enumerate(chunks):
                embedding_vec: list[float] | None = None
                try:
                    raw_embedding = await client.embed(chunk_text, user_id=user_id)
                    if raw_embedding is not None:
                        embedding_vec = raw_embedding if isinstance(raw_embedding, list) else list(raw_embedding)
                except Exception as emb_exc:
                    log.warning(
                        'kb_chunk_embed_failed',
                        document_id=document_id,
                        chunk_index=idx,
                        error=str(emb_exc),
                    )

                token_count = _estimate_tokens(chunk_text)
                await conn.execute(
                    """
                    INSERT INTO kb_chunks (document_id, user_id, chunk_index, content, embedding, token_count)
                    VALUES ($1, $2, $3, $4, $5, $6)
                    """,
                    document_id,
                    user_id,
                    idx,
                    chunk_text,
                    np.array(embedding_vec, dtype=np.float32) if embedding_vec else None,
                    token_count,
                )
                inserted += 1

            await conn.execute(
                """
                UPDATE kb_documents
                SET status = 'ready', chunk_count = $1, error_message = NULL, updated_at = NOW()
                WHERE id = $2
                """,
                inserted,
                document_id,
            )

        log.info(
            'kb_document_processing_done',
            document_id=document_id,
            chunks_inserted=inserted,
        )

    except Exception as exc:
        log.error('kb_document_processing_failed', document_id=document_id, error=str(exc))
        async with pool.acquire() as conn:
            await conn.execute(
                """
                UPDATE kb_documents
                SET status = 'error', error_message = $1, updated_at = NOW()
                WHERE id = $2
                """,
                str(exc)[:500],
                document_id,
            )


async def retrieve_relevant_chunks(
    user_id: str,
    agent_id: str | None,
    query: str,
    limit: int = 5,
) -> list[dict]:
    """
    Retrieve the most semantically relevant KB chunks for a query.

    Generates a query embedding, then performs a cosine-distance vector search
    against kb_chunks for the user. Agent lookups include agent-specific docs,
    global docs, and general company/business docs.

    Returns a list of dicts with 'content' and 'score' keys, sorted by
    relevance descending. Returns an empty list if embeddings are unavailable
    or no chunks exist.
    """
    if not query.strip():
        return []

    client = get_ai_client()

    try:
        query_embedding = await client.embed(query[:2000], user_id=user_id)
    except Exception as exc:
        log.warning('kb_query_embed_failed', error=str(exc))
        query_embedding = None

    if query_embedding is None:
        return await _keyword_retrieve(user_id, agent_id, query, limit)

    query_vec = np.array(
        query_embedding if isinstance(query_embedding, list) else list(query_embedding),
        dtype=np.float32,
    )

    pool = await get_pool()
    async with pool.acquire() as conn:
        if agent_id:
            # Include chunks from this agent's documents and the global KB (agent_id IS NULL)
            rows = await conn.fetch(
                """
                SELECT kc.content, 1 - (kc.embedding <-> $1) AS score,
                       kd.id AS document_id, kd.title AS document_title,
                       kd.source_type, kd.category
                FROM kb_chunks kc
                JOIN kb_documents kd ON kd.id = kc.document_id
                WHERE kc.user_id = $2
                  AND kc.embedding IS NOT NULL
                  AND (
                    kd.agent_id = $3
                    OR kd.agent_id IS NULL
                    OR kd.category ILIKE ANY(ARRAY['company','business','policies','products','pricing'])
                  )
                  AND kd.status = 'ready'
                ORDER BY kc.embedding <-> $1
                LIMIT $4
                """,
                query_vec,
                user_id,
                agent_id,
                limit,
            )
        else:
            rows = await conn.fetch(
                """
                SELECT kc.content, 1 - (kc.embedding <-> $1) AS score,
                       kd.id AS document_id, kd.title AS document_title,
                       kd.source_type, kd.category
                FROM kb_chunks kc
                JOIN kb_documents kd ON kd.id = kc.document_id
                WHERE kc.user_id = $2
                  AND kc.embedding IS NOT NULL
                  AND kd.status = 'ready'
                ORDER BY kc.embedding <-> $1
                LIMIT $3
                """,
                query_vec,
                user_id,
                limit,
            )

        if rows:
            doc_ids = list({row['document_id'] for row in rows})
            await conn.execute(
                'UPDATE kb_documents SET used_count = used_count + 1, last_used_at = NOW() WHERE id = ANY($1::uuid[])',
                doc_ids,
            )

    results = [
        {
            'content': row['content'],
            'score': float(row['score']),
            'document_id': str(row['document_id']),
            'document_title': row['document_title'],
            'source_type': row['source_type'],
            'category': row['category'],
        }
        for row in rows
    ]

    log.debug(
        'kb_chunks_retrieved',
        user_id=user_id,
        agent_id=agent_id,
        query_length=len(query),
        results_found=len(results),
    )

    return results


async def _keyword_retrieve(user_id: str, agent_id: str | None, query: str, limit: int) -> list[dict]:
    terms = [term for term in re.findall(r'[A-Za-z0-9]{3,}', query.lower())[:8]]
    if not terms:
        return []
    pool = await get_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch(
            """
            SELECT kc.content,
                   0.35::float AS score,
                   kd.id AS document_id,
                   kd.title AS document_title,
                   kd.source_type,
                   kd.category
            FROM kb_chunks kc
            JOIN kb_documents kd ON kd.id = kc.document_id
            WHERE kc.user_id = $1
              AND kd.status = 'ready'
              AND ($2::uuid IS NULL OR kd.agent_id = $2 OR kd.agent_id IS NULL)
              AND EXISTS (
                SELECT 1 FROM unnest($3::text[]) term
                WHERE kc.content ILIKE '%' || term || '%'
              )
            ORDER BY kd.updated_at DESC, kc.chunk_index ASC
            LIMIT $4
            """,
            user_id,
            agent_id,
            terms,
            limit,
        )

    return [
        {
            'content': row['content'],
            'score': float(row['score']),
            'document_id': str(row['document_id']),
            'document_title': row['document_title'],
            'source_type': row['source_type'],
            'category': row['category'],
        }
        for row in rows
    ]


async def search_knowledge(user_id: str, query: str, limit: int = 10) -> list[dict]:
    """Search the knowledge base and return matching chunks with document metadata."""
    if not query.strip():
        return []

    from ..ai.client import get_ai_client
    client = get_ai_client()

    try:
        query_embedding = await client.embed(query[:2000], user_id=user_id)
    except Exception as exc:
        log.warning('kb_search_embed_failed', error=str(exc))
        return []

    if query_embedding is None:
        return []

    query_vec = __import__('numpy').array(
        query_embedding if isinstance(query_embedding, list) else list(query_embedding),
        dtype=__import__('numpy').float32,
    )

    pool = await get_pool()
    async with pool.acquire() as conn:
        rows = await conn.fetch(
            """
            SELECT kc.content,
                   1 - (kc.embedding <-> $1) AS score,
                   kd.id AS document_id,
                   kd.title AS document_title,
                   kd.source_type,
                   kd.category
            FROM kb_chunks kc
            JOIN kb_documents kd ON kd.id = kc.document_id
            WHERE kc.user_id = $2
              AND kc.embedding IS NOT NULL
              AND kd.status = 'ready'
            ORDER BY kc.embedding <-> $1
            LIMIT $3
            """,
            query_vec,
            user_id,
            limit,
        )
        if rows:
            doc_ids = list({row['document_id'] for row in rows})
            for doc_id in doc_ids:
                await conn.execute(
                    "UPDATE kb_documents SET used_count = used_count + 1, last_used_at = NOW() WHERE id = $1",
                    doc_id,
                )

    return [
        {
            'content': row['content'],
            'score': float(row['score']),
            'document_id': str(row['document_id']),
            'document_title': row['document_title'],
            'source_type': row['source_type'],
            'category': row['category'],
        }
        for row in rows
    ]


async def chat_with_knowledge(user_id: str, question: str) -> dict:
    """Answer a question using the knowledge base as context."""
    chunks = await retrieve_relevant_chunks(user_id, None, question, limit=5)

    if not chunks:
        return {
            'answer': "I don't have information about that in my knowledge base. Try adding relevant documents first.",
            'sources': [],
        }

    context_parts = [f"[Source {i+1}]: {c['content']}" for i, c in enumerate(chunks)]
    context = '\n\n'.join(context_parts)

    prompt = f"""You are an AI assistant answering questions from a business knowledge base.

Question: {question}

Knowledge Base Content:
{context}

Answer the question accurately and concisely based only on the knowledge base content above.
If the information is not in the knowledge base, say so clearly.
Do not make up information."""

    from ..ai.client import get_ai_client
    client = get_ai_client()
    try:
        answer = await client.complete_text(
            [{'role': 'user', 'content': prompt}],
            service='intelligence', feature='knowledge_qa', user_id=user_id,
        )
    except Exception as exc:
        log.error('kb_chat_failed', error=str(exc))
        answer = 'Sorry, I was unable to generate an answer. Please try again.'

    sources = [
        {
            'content': c['content'][:300] + ('...' if len(c['content']) > 300 else ''),
            'score': round(c['score'], 3),
        }
        for c in chunks[:3]
    ]

    return {'answer': answer or '', 'sources': sources}
